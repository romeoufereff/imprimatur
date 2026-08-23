#!/usr/bin/env python3
"""
probe_brand.py — read a brand artefact and record what is ACTUALLY in it.

Emits one evidence.json: colours with area weights and where they were seen, type sizes with
their fonts and sample strings, font families by frequency, and any gradients readable from the
source. Nothing here decides what a brand's tokens ARE — that judgment belongs to the skill.
This script's only job is to make the judgment checkable, by attaching a count and an area to
every value so a human can see why a colour is in the pack.

Why area weighting: frequency alone ranks a hairline rule above a full-bleed cover. A brand's
structural colours are the ones that cover the page, so every observation carries the fraction
of the page it painted, and colours are ranked by summed area.

Why perceptual merging: brand decks are full of near-duplicates — antialiasing, JPEG artefacts,
a rectangle nudged to #0048FE by hand. Merging in Lab space at a ΔE threshold collapses those
into one token instead of minting three. Where a declared theme exists (.pptx), sampled colours
SNAP to the theme's exact hex when they are within threshold, so the pack carries the designer's
value rather than a rendering of it.

Usage:
    probe_brand.py <file.pptx|.potx|.pdf> [--out evidence.json]
                   [--pages 1-30] [--delta-e 3.0] [--min-area 0.002] [--raster]

    --raster   PDF only: also rasterise pages and sample dominant colours. Needed when the PDF
               is a flattened export (few or no vector objects); slower, and its colours are
               marked "raster" because they are pixels, not declared values.
"""
import argparse
import json
import os
import re
import sys
import zipfile
from collections import defaultdict
from xml.etree import ElementTree as ET

A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


# ── colour maths ───────────────────────────────────────────────────────────────
def _srgb_to_lab(hex_str):
    """CIE L*a*b* under D65. Perceptual distance is what makes 'near-duplicate' meaningful."""
    r, g, b = (int(hex_str[i:i + 2], 16) / 255.0 for i in (1, 3, 5))

    def lin(c):
        return c / 12.92 if c <= 0.04045 else ((c + 0.055) / 1.055) ** 2.4

    r, g, b = lin(r), lin(g), lin(b)
    x = (r * 0.4124 + g * 0.3576 + b * 0.1805) / 0.95047
    y = (r * 0.2126 + g * 0.7152 + b * 0.0722) / 1.00000
    z = (r * 0.0193 + g * 0.1192 + b * 0.9505) / 1.08883

    def f(t):
        return t ** (1 / 3) if t > 0.008856 else (7.787 * t) + (16 / 116)

    fx, fy, fz = f(x), f(y), f(z)
    return (116 * fy - 16, 500 * (fx - fy), 200 * (fy - fz))


def _delta_e(a, b):
    return sum((x - y) ** 2 for x, y in zip(a, b)) ** 0.5


def norm_font(name):
    """Strip PDF subset prefixes so one family stops looking like twenty.

    Embedded PDF fonts are named `AAAAAG+Calibri` — a per-subset tag plus the real family.
    Left alone, a deck using two weights reports dozens of 'fonts' and buries the fact that
    the brand has one typeface. The tag carries no design information, so it goes.
    """
    if not name:
        return name
    return re.sub(r"^[A-Z]{6}\+", "", str(name)).strip()


def _hex(rgb):
    return "#%02X%02X%02X" % tuple(max(0, min(255, int(round(c)))) for c in rgb)


def _rel_luminance(hex_str):
    def lin(c):
        c /= 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    r, g, b = (int(hex_str[i:i + 2], 16) for i in (1, 3, 5))
    return 0.2126 * lin(r) + 0.7152 * lin(g) + 0.0722 * lin(b)


def contrast_ratio(fg, bg):
    a, b = _rel_luminance(fg), _rel_luminance(bg)
    hi, lo = max(a, b), min(a, b)
    return round((hi + 0.05) / (lo + 0.05), 2)


# ── observation store ──────────────────────────────────────────────────────────
class Evidence:
    def __init__(self):
        self.colors = []   # {hex, area, kind, page, note}
        self.type = []     # {px, font, page, sample}
        self.fonts = defaultdict(int)
        self.gradients = []
        self.notes = []
        self.pages_seen = set()

    def color(self, hex_str, area, kind, page, note=""):
        if not hex_str:
            return
        # Clamp to one page: bleed shapes and grouped duplicates routinely report more than
        # the slide they sit on, and an uncapped outlier would outrank the real structure.
        self.colors.append({"hex": hex_str.upper(), "area": min(1.0, max(0.0, float(area or 0.0))),
                            "kind": kind, "page": page, "note": note})
        self.pages_seen.add(page)

    def text(self, px, font, page, sample):
        font = norm_font(font)
        if px and px > 0:
            self.type.append({"px": round(float(px), 1), "font": font or "?",
                              "page": page, "sample": (sample or "").strip()[:60]})
            if font:
                self.fonts[font] += 1

    def cluster_colors(self, delta_e, min_area, snap_to=None):
        """Merge near-duplicates, ranked by summed area. `snap_to` = declared theme hexes."""
        buckets = []
        for obs in sorted(self.colors, key=lambda o: -o["area"]):
            lab = _srgb_to_lab(obs["hex"])
            for b in buckets:
                if _delta_e(lab, b["lab"]) <= delta_e:
                    b["members"].append(obs)
                    break
            else:
                buckets.append({"lab": lab, "rep": obs["hex"], "members": [obs]})

        declared = {h.upper(): _srgb_to_lab(h) for h in (snap_to or [])}
        out = []
        # area_weight is the mean page-areas this colour PAINTS per page. It is a weight, not
        # a share: overlapping objects paint the same pixels repeatedly, so a background under
        # a dozen white boxes legitimately exceeds 1.0. That is fine for ranking — which is all
        # it is for — but do not read it as "percent of the deck".
        npages = max(1, len(self.pages_seen))
        for b in buckets:
            area = sum(m["area"] for m in b["members"]) / npages
            rep, snapped = b["rep"], None
            for dhex, dlab in declared.items():
                if _delta_e(b["lab"], dlab) <= delta_e:
                    rep, snapped = dhex, dhex
                    break
            kinds = defaultdict(int)
            pages = set()
            for m in b["members"]:
                kinds[m["kind"]] += 1
                pages.add(m["page"])
            out.append({
                "hex": rep, "snapped_to_theme": snapped,
                "area_weight": round(area, 5), "count": len(b["members"]),
                "pages": len(pages), "kinds": dict(kinds),
                "variants": sorted({m["hex"] for m in b["members"]} - {rep})[:6],
            })
        out.sort(key=lambda c: (-c["area_weight"], -c["count"]))
        # A colour the brand DECLARED survives regardless of how sparingly the file uses it:
        # the theme is the designer's statement of intent, and an accent that appears on two
        # pages is still a brand token. Usage thresholds only arbitrate undeclared colours,
        # which is where the antialiasing noise lives.
        kept = [c for c in out
                if c["snapped_to_theme"] or c["area_weight"] >= min_area or c["pages"] >= 3]
        return kept, len(out) - len(kept)

    def type_ramp(self, tol=0.6, min_count=2):
        steps = []
        for t in sorted(self.type, key=lambda t: -t["px"]):
            for s in steps:
                if abs(s["px"] - t["px"]) <= tol:
                    s["members"].append(t)
                    break
            else:
                steps.append({"px": t["px"], "members": [t]})
        out = []
        for s in steps:
            fonts = defaultdict(int)
            pages = set()
            for m in s["members"]:
                fonts[m["font"]] += 1
                pages.add(m["page"])
            out.append({
                "px": s["px"], "count": len(s["members"]), "pages": len(pages),
                "fonts": dict(sorted(fonts.items(), key=lambda kv: -kv[1])[:4]),
                "samples": [m["sample"] for m in s["members"][:3] if m["sample"]],
            })
        return [s for s in out if s["count"] >= min_count or s["pages"] >= 2]


# ── PPTX ───────────────────────────────────────────────────────────────────────
def read_theme(path):
    """The declared design system. A .pptx carries twelve named colours and a font pair that
    the brand's own designer chose — ground truth, not inference. Read straight from the XML so
    a corrupt or unusual deck degrades to 'no theme' instead of raising."""
    theme = {"colors": {}, "fonts": {}}
    try:
        with zipfile.ZipFile(path) as z:
            names = [n for n in z.namelist() if re.match(r"ppt/theme/theme\d+\.xml$", n)]
            if not names:
                return theme
            root = ET.fromstring(z.read(sorted(names)[0]))
        scheme = root.find(f".//{A_NS}clrScheme")
        if scheme is not None:
            for child in scheme:
                slot = child.tag.split("}")[-1]
                srgb = child.find(f"{A_NS}srgbClr")
                sysc = child.find(f"{A_NS}sysClr")
                if srgb is not None:
                    theme["colors"][slot] = "#" + srgb.get("val", "").upper()
                elif sysc is not None and sysc.get("lastClr"):
                    theme["colors"][slot] = "#" + sysc.get("lastClr").upper()
        fonts = root.find(f".//{A_NS}fontScheme")
        if fonts is not None:
            for slot, tag in (("major", "majorFont"), ("minor", "minorFont")):
                node = fonts.find(f"{A_NS}{tag}/{A_NS}latin")
                if node is not None and node.get("typeface"):
                    theme["fonts"][slot] = node.get("typeface")
    except (zipfile.BadZipFile, ET.ParseError, KeyError, OSError):
        pass
    return theme


def probe_pptx(path, ev, max_pages):
    from pptx import Presentation
    from pptx.util import Emu

    theme = read_theme(path)
    scheme_hex = dict(theme["colors"])
    # python-pptx names theme slots differently from the XML; map both spellings.
    alias = {"ACCENT_1": "accent1", "ACCENT_2": "accent2", "ACCENT_3": "accent3",
             "ACCENT_4": "accent4", "ACCENT_5": "accent5", "ACCENT_6": "accent6",
             "DARK_1": "dk1", "DARK_2": "dk2", "LIGHT_1": "lt1", "LIGHT_2": "lt2",
             "HYPERLINK": "hlink", "FOLLOWED_HYPERLINK": "folHlink", "TEXT_1": "dk1",
             "TEXT_2": "dk2", "BACKGROUND_1": "lt1", "BACKGROUND_2": "lt2"}

    prs = Presentation(path)
    slide_area = float(Emu(prs.slide_width) * Emu(prs.slide_height)) or 1.0

    def resolve(color_fmt):
        try:
            if color_fmt is None or color_fmt.type is None:
                return None
            if str(color_fmt.type).startswith("MSO_THEME_COLOR") or hasattr(color_fmt, "theme_color"):
                try:
                    name = str(color_fmt.theme_color).split(".")[-1].split(" ")[0]
                    if name in alias and alias[name] in scheme_hex:
                        return scheme_hex[alias[name]]
                except (AttributeError, ValueError):
                    pass
            return "#" + str(color_fmt.rgb)
        except (AttributeError, TypeError, ValueError):
            return None

    for i, slide in enumerate(prs.slides, start=1):
        if max_pages and i > max_pages:
            break
        for shape in slide.shapes:
            try:
                area = 0.0
                if shape.width and shape.height:
                    area = float(shape.width) * float(shape.height) / slide_area
                fill = getattr(shape, "fill", None)
                if fill is not None:
                    ftype = str(getattr(fill, "type", ""))
                    if "GRADIENT" in ftype:
                        stops = []
                        for st in fill.gradient_stops:
                            h = resolve(st.color)
                            if h:
                                stops.append({"pos": round(float(st.position), 3), "hex": h})
                                ev.color(h, area / max(1, len(fill.gradient_stops)),
                                         "gradient", i, shape.shape_type and str(shape.shape_type))
                        if len(stops) >= 2:
                            ev.gradients.append({"page": i, "stops": stops,
                                                 "angle": getattr(fill, "gradient_angle", None)})
                    elif "SOLID" in ftype:
                        ev.color(resolve(fill.fore_color), area, "fill", i)
                line = getattr(shape, "line", None)
                if line is not None and getattr(line, "color", None) is not None:
                    ev.color(resolve(line.color), area * 0.02, "line", i)
            except (AttributeError, TypeError, ValueError, KeyError):
                pass

            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        px = float(run.font.size.pt) * (96.0 / 72.0) if run.font.size else None
                        name = run.font.name or theme["fonts"].get("minor")
                        if run.font.bold:
                            name = f"{name} Bold" if name else "Bold"
                        ev.text(px, name, i, run.text)
                        h = resolve(run.font.color)
                        if h and run.text.strip():
                            # Text area is glyph-ish, not the shape box: weight it down so a
                            # single word cannot outrank a background.
                            ev.color(h, area * 0.05, "text", i)
                    except (AttributeError, TypeError, ValueError):
                        pass

    if theme["colors"]:
        ev.notes.append(f"declared theme found: {len(theme['colors'])} colours, "
                        f"fonts {theme['fonts'] or 'none'}")
    else:
        ev.notes.append("no theme1.xml colour scheme — treat every colour as inferred")
    return theme


# ── PDF ────────────────────────────────────────────────────────────────────────
def _pdf_color(val):
    if val is None:
        return None
    try:
        seq = list(val) if isinstance(val, (list, tuple)) else [float(val)]
    except (TypeError, ValueError):
        return None
    try:
        if len(seq) == 1:
            g = float(seq[0]) * 255
            return _hex((g, g, g))
        if len(seq) == 3:
            return _hex(tuple(float(c) * 255 for c in seq))
        if len(seq) == 4:
            c, m, y, k = (float(x) for x in seq)
            return _hex((255 * (1 - c) * (1 - k), 255 * (1 - m) * (1 - k), 255 * (1 - y) * (1 - k)))
    except (TypeError, ValueError):
        return None
    return None


def probe_pdf(path, ev, max_pages, raster):
    import pdfplumber

    with pdfplumber.open(path) as pdf:
        pages = pdf.pages[:max_pages] if max_pages else pdf.pages
        for i, page in enumerate(pages, start=1):
            page_area = float(page.width * page.height) or 1.0
            for kind, items in (("fill", page.rects), ("fill", page.curves)):
                for it in items:
                    try:
                        h = _pdf_color(it.get("non_stroking_color"))
                        w = abs(float(it.get("x1", 0)) - float(it.get("x0", 0)))
                        ht = abs(float(it.get("bottom", 0)) - float(it.get("top", 0)))
                        ev.color(h, (w * ht) / page_area, kind, i)
                    except (TypeError, ValueError):
                        pass
            words = defaultdict(list)
            for ch in page.chars:
                try:
                    key = (round(float(ch.get("size", 0)), 1), ch.get("fontname", "?"))
                    words[key].append(ch)
                    h = _pdf_color(ch.get("non_stroking_color"))
                    w = abs(float(ch.get("x1", 0)) - float(ch.get("x0", 0)))
                    ht = abs(float(ch.get("bottom", 0)) - float(ch.get("top", 0)))
                    ev.color(h, (w * ht) / page_area, "text", i)
                except (TypeError, ValueError):
                    pass
            for (size, font), chars in words.items():
                # PDF sizes are points in user space; the pack works in CSS px.
                ev.text(size * (96.0 / 72.0), font, i,
                        "".join(c.get("text", "") for c in chars[:60]))

    if raster:
        _probe_raster(path, ev, max_pages)
    ev.notes.append("PDF carries no declared theme — every value here is inferred from what "
                    "the pages happen to contain")
    if not ev.gradients:
        # Measured on a PPTX and its PowerPoint PDF export: 12 gradients readable from the
        # source, 0 from the PDF. Export flattens them into shading patterns that carry no
        # stop list, so a brand gradient must come from the source file or be reconstructed
        # by hand from a screenshot. Silence here would read as "the brand has no gradient".
        ev.notes.append("no gradients readable — PDF export flattens gradients into shading "
                        "patterns with no stop list. If the brand has one, get the .pptx or "
                        "rebuild it by hand; do not conclude the brand has none")


def _probe_raster(path, ev, max_pages):
    """Flattened exports have no vector objects to read; sample the pixels instead."""
    try:
        import pypdfium2 as pdfium
        import numpy as np
    except ImportError:
        ev.notes.append("--raster requested but pypdfium2/numpy unavailable; skipped")
        return
    try:
        doc = pdfium.PdfDocument(path)
    except Exception as e:  # noqa: BLE001 - pdfium raises a variety of load errors
        ev.notes.append(f"--raster failed to open the PDF ({e}); skipped")
        return
    n = min(len(doc), max_pages or len(doc))
    for i in range(n):
        try:
            arr = np.asarray(doc[i].render(scale=0.6).to_pil().convert("RGB"))
        except Exception:  # noqa: BLE001 - a single unrenderable page must not abort the probe
            continue
        flat = (arr // 8 * 8).reshape(-1, 3)
        colors, counts = np.unique(flat, axis=0, return_counts=True)
        total = counts.sum() or 1
        order = np.argsort(-counts)[:12]
        for idx in order:
            share = float(counts[idx]) / total
            if share < 0.005:
                continue
            ev.color(_hex(tuple(colors[idx])), share, "raster", i + 1)
    ev.notes.append(f"rastered {n} page(s); raster colours are pixels, not declared values")


# ── main ───────────────────────────────────────────────────────────────────────
def parse_pages(spec):
    if not spec:
        return None
    m = re.match(r"^(?:(\d+)-)?(\d+)$", spec.strip())
    return int(m.group(2)) if m else None


def main():
    ap = argparse.ArgumentParser(description="Record the tokens actually present in a brand artefact.")
    ap.add_argument("source")
    ap.add_argument("--out", default="evidence.json")
    ap.add_argument("--pages", default=None, help="Cap pages, e.g. 30 or 1-30")
    ap.add_argument("--delta-e", type=float, default=3.0,
                    help="Perceptual merge threshold; 3.0 is roughly 'a trained eye can just tell'")
    ap.add_argument("--min-area", type=float, default=0.002)
    ap.add_argument("--raster", action="store_true", help="PDF: also sample rasterised pages")
    args = ap.parse_args()

    src = os.path.abspath(os.path.expanduser(args.source))
    if not os.path.isfile(src):
        sys.exit(f"error: no such file: {src}")
    ext = os.path.splitext(src)[1].lower()
    max_pages = parse_pages(args.pages)

    ev = Evidence()
    theme = {"colors": {}, "fonts": {}}
    if ext in (".pptx", ".potx"):
        theme = probe_pptx(src, ev, max_pages)
        kind = "pptx"
    elif ext == ".pdf":
        probe_pdf(src, ev, max_pages, args.raster)
        kind = "pdf"
    else:
        sys.exit(f"error: unsupported source {ext} — this skill reads .pptx, .potx or .pdf")

    colors, dropped = ev.cluster_colors(args.delta_e, args.min_area,
                                        snap_to=list(theme["colors"].values()))
    ramp = ev.type_ramp()
    out = {
        "source": {"path": src, "kind": kind, "observations": len(ev.colors)},
        "theme": theme,
        "colors": colors,
        "dropped_as_noise": dropped,
        "type": ramp,
        "fonts": dict(sorted(ev.fonts.items(), key=lambda kv: -kv[1])[:20]),
        "gradients": ev.gradients[:12],
        "notes": ev.notes,
    }
    with open(args.out, "w", encoding="utf-8") as f:
        json.dump(out, f, indent=2, ensure_ascii=False)
        f.write("\n")

    print(f"{kind}: {len(ev.colors)} colour observations -> {len(colors)} clusters "
          f"({dropped} dropped as noise), {len(ramp)} type steps, {len(ev.fonts)} font(s)")
    for c in colors[:8]:
        snap = "  [theme]" if c["snapped_to_theme"] else ""
        print(f"  {c['hex']}  weight {c['area_weight']:.4f}  {c['pages']} page(s)  "
              f"{c['kinds']}{snap}")
    for n in ev.notes:
        print(f"  note: {n}")
    print(f"-> {args.out}")


if __name__ == "__main__":
    main()
