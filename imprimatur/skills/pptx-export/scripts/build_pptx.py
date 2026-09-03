#!/usr/bin/env python3
"""
pptx-export — stage 2: IR JSON → PowerPoint.

Renders the IR emitted by extract_ir.py into a 16:9 .pptx via python-pptx:

  slide bg → solid fill, or the full-slide screenshot as a base picture when the
             background is a gradient (covers, section dividers)
  box      → (rounded) rectangle, solid or gradient fill / border, shadows off
  text     → textbox with styled runs; px → pt at 0.5 (540pt slide / 1080px canvas)
  raster   → the element screenshot placed as a picture

Geometry: 1920x1080 px maps onto a 13.333in x 7.5in slide at exactly 6350 EMU/px.

Usage (html2pptx.py drives the whole deck):
    python3 build_pptx.py ir/01.json ir/02.json ... --output deck.pptx
"""

import argparse
import json
import math
import os
import re
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

EMU_PER_PX = 6350          # 12192000 EMU / 1920 px
# 1080 px map onto a 7.5in-tall slide (540 pt) -> 1 px = 0.5 pt. Using the 96-dpi
# CSS convention (x0.75) here renders ALL text 1.5x too large — the confirmed root
# cause of the 'far from HTML' PowerPoint look. Never reintroduce it.
PT_PER_PX = 0.5
SLIDE_W = Emu(12192000)    # 13.333 in
SLIDE_H = Emu(6858000)     # 7.5 in

ALIGN = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER,
         'right': PP_ALIGN.RIGHT, 'justify': PP_ALIGN.JUSTIFY}

A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
SVG_EXT_URI = '{96DAC541-7B7A-4930-B5D3-8C156C88E4B4}'


# Brand values (fonts, chart colors) come from the ACTIVE design system, never from
# constants in this file — that is what lets the same exporter serve any pack.
# Depth-independent: anchor on the plugin marker rather than counting `..` hops,
# so moving this file between skill folders cannot silently resolve to the wrong
# directory. See scripts/_paths.py.
_here = os.path.dirname(os.path.abspath(__file__))
_probe = _here
while not os.path.isdir(os.path.join(_probe, ".claude-plugin")):
    _parent = os.path.dirname(_probe)
    if _parent == _probe:
        raise RuntimeError(f"No .claude-plugin/ found above {_here}")
    _probe = _parent
sys.path.insert(0, os.path.join(_probe, "scripts"))
from ds_config import load as _load_ds  # noqa: E402

_DS = _load_ds()
_FONT_BY_WEIGHT = sorted(
    ((int(k), v) for k, v in (_DS.get("pptx.fontFamilyByWeight", {}) or {}).items()),
    reverse=True)


def font_for_weight(w):
    """Resolve a numeric weight onto the pack's installed per-weight family name.

    PowerPoint substitutes by family NAME: writing the bare family plus a bold flag
    collapses everything to one mid weight (the v1 'everything looks semi-bold' bug),
    so the pack declares the exact full name per weight band.
    """
    for threshold, name in _FONT_BY_WEIGHT:
        if w >= threshold:
            return name
    return _FONT_BY_WEIGHT[-1][1] if _FONT_BY_WEIGHT else "sans-serif"


def px(v):
    return Emu(int(round(v * EMU_PER_PX)))


def rgb(hexstr):
    return RGBColor.from_string(hexstr.lstrip('#'))


def add_box(slide, node):
    r = node['rect']
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if node.get('radius', 0) > 0 else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shape_type, px(r['x']), px(r['y']), px(r['w']), px(r['h']))
    if shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        # adjustment 0..0.5 = radius as fraction of the shorter side
        frac = min(0.5, node['radius'] / max(1.0, min(r['w'], r['h'])))
        try:
            sp.adjustments[0] = frac
        except Exception:
            pass
    if node.get('gradient'):
        # Brand-ramp badge/dot/bar: a fully-opaque element whose background-image was
        # exactly one simple linear-gradient (extract_ir.py's elGrad) — same native
        # gradFill path as the slide background / gradient-text runs, not a picture.
        grad = node['gradient']
        sp.fill.solid()  # creates the fill element in schema position, then replaced
        solid = sp._element.spPr.find(qn('a:solidFill'))
        stops = ''.join(
            f'<a:gs pos="{int(round(s["pos"] * 1000))}"><a:srgbClr val="{s["color"].lstrip("#")}"/></a:gs>'
            for s in grad['stops'])
        # CSS Ndeg (0=up, cw) -> DrawingML (0=right, cw, 1/60000 deg): (N-90)%360
        ang = int(round(((grad['angle'] - 90) % 360) * 60000))
        gfill = parse_xml(
            f'<a:gradFill xmlns:a="{A_NS}" rotWithShape="1"><a:gsLst>{stops}</a:gsLst>'
            f'<a:lin ang="{ang}" scaled="1"/></a:gradFill>')
        solid.getparent().replace(solid, gfill)
    elif node.get('fill'):
        sp.fill.solid()
        # semi-transparent fills are pre-composited over white — a v1 approximation
        # (true per-shape alpha needs raw oxml; not worth it for the rare cases)
        a = node.get('fillAlpha', 1)
        c = rgb(node['fill'])
        if a < 1:
            c = RGBColor(*(int(round(v * a + 255 * (1 - a))) for v in (c[0], c[1], c[2])))
        sp.fill.fore_color.rgb = c
    else:
        sp.fill.background()
    if node.get('border'):
        sp.line.color.rgb = rgb(node['border']['color'])
        sp.line.width = Pt(node['border']['width'] * PT_PER_PX)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def _apply_para_props(para, node, para_runs):
    para.alignment = ALIGN.get(node.get('align', 'left'), PP_ALIGN.LEFT)
    # exact line spacing from computed px line-height — per PARAGRAPH, from its own
    # runs (a card can mix a 28px title line with 20px body lines; the container's
    # single line-height would squeeze the title). Falls back to the node value.
    lhs = [r['st'].get('lineHeight') for r in para_runs if r.get('st') and r['st'].get('lineHeight')]
    lh = max(lhs) if lhs else node.get('lineHeight')
    if lh:
        para.line_spacing = Pt(round(lh * PT_PER_PX, 1))


def add_text(slide, node):
    r = node['rect']
    # small breathing room so wrapped lines match the HTML box
    tb = slide.shapes.add_textbox(px(r['x'] - 1), px(r['y'] - 2), px(r['w'] + 4), px(r['h'] + 6))
    tf = tb.text_frame
    # Wrapping is ALWAYS off: the extractor measured the browser's actual line breaks
    # and emitted them as explicit '\n' runs, so every renderer reproduces the same
    # lines. With wrapping on, a renderer whose font metrics differ by 2% reflows a
    # title into the subtitle below it — the v1 failure mode.
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    # kill autofit so PPT never rescales the text to the box
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    for fit in ('a:normAutofit', 'a:spAutoFit'):
        el = bodyPr.find(qn(fit))
        if el is not None:
            bodyPr.remove(el)
    # group runs into paragraphs first so each paragraph gets its own line-height
    paragraphs, current = [], []
    for run in node['runs']:
        if run['text'] == '\n':
            paragraphs.append(current)
            current = []
        else:
            current.append(run)
    paragraphs.append(current)

    for i, para_runs in enumerate(paragraphs):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _apply_para_props(para, node, para_runs)
        for run in para_runs:
            rn = para.add_run()
            rn.text = run['text']
            st = run['st']
            rn.font.name = font_for_weight(st['weight'])
            rn.font.size = Pt(round(st['size'] * PT_PER_PX, 1))
            rn.font.bold = False      # weight is carried by the per-weight font name
            rn.font.italic = bool(st.get('italic'))
            rn.font.color.rgb = rgb(st['color'])
            if st.get('accent'):
                _set_run_brand_gradient(rn, st.get('accentGradient'))
            if st.get('spacing'):
                # letter-spacing: a:rPr spc is in 1/100 pt
                rn.font._rPr.set('spc', str(int(round(st['spacing'] * PT_PER_PX * 100))))
    return tb


def _parse_css_gradient(css):
    """'linear-gradient(60deg, #2B59C3 0%, #3F7FD1 55%, #6E4BB8 100%)' ->
    (angle_deg, [(pos_fraction, '#hex'), ...]), or None if the shape doesn't match
    (radial, no explicit stop percentages, etc.) — the pack declares whatever it wants;
    only the shape this engine knows how to render is usable here."""
    m = re.match(r'linear-gradient\(\s*([\d.]+)deg\s*,\s*(.+)\)\s*$', css.strip())
    if not m:
        return None
    stops = []
    for part in m.group(2).split(','):
        sm = re.match(r'\s*(#[0-9A-Fa-f]{6})\s+([\d.]+)%\s*$', part.strip())
        if not sm:
            return None
        stops.append((float(sm.group(2)) / 100.0, sm.group(1)))
    return float(m.group(1)), stops


def _set_run_brand_gradient(rn, css_override=None):
    """Gradient-text runs (background-clip:text in HTML) get a REAL per-run gradient
    fill — read from the DECK'S OWN computed background-image where the IR carried one
    (css_override — respects a deck's !important brand-ramp override, e.g. dropping
    violet), falling back to the active pack's generic 'brand' gradient only when the
    IR has no per-run value (older IR, or an unparseable override). Never hardcoded.
    PowerPoint renders run-level gradFill natively. The gradient spans each run's own
    box (HTML spans the whole element) — a close approximation for accent phrases."""
    rPr = rn.font._rPr
    solid = rPr.find(qn('a:solidFill'))
    if solid is None:
        return
    css = css_override or _DS.gradients().get('brand')
    parsed = _parse_css_gradient(css) if css else None
    if parsed is None and css_override:
        # deck override didn't parse (e.g. layered/radial) — fall back to the pack's own
        css = _DS.gradients().get('brand')
        parsed = _parse_css_gradient(css) if css else None
    if parsed is None:
        return  # pack has no usable 'brand' gradient — leave the flat colour already set
    angle, stops = parsed
    gs = ''.join(
        f'<a:gs pos="{int(round(p * 100000))}"><a:srgbClr val="{c.lstrip("#")}"/></a:gs>'
        for p, c in stops)
    # CSS Ndeg (0=up, cw) -> DrawingML (0=right, cw, 1/60000 deg): (N-90)%360
    ang = int(round(((angle - 90) % 360) * 60000))
    grad = parse_xml(
        f'<a:gradFill xmlns:a="{A_NS}" rotWithShape="1"><a:gsLst>{gs}</a:gsLst>'
        f'<a:lin ang="{ang}" scaled="1"/></a:gradFill>')
    solid.getparent().replace(solid, grad)


def add_raster(slide, node, svg_mode='shapes'):
    """Placement of non-text visuals. SVGs become NATIVE GROUPED SHAPES by default
    (svg2shapes: rects/ovals/freeform beziers/textboxes — movable, scalable, editable,
    ungroupable in PowerPoint). Unconvertible SVGs (curved textPath labels, masks,
    chart-sized node counts) fall back to svgBlip: vector picture + PNG fallback,
    manually convertible via right-click → Convert to Shape. `--svg-blip` forces the
    blip path for everything."""
    if not node.get('asset'):
        return None
    r = node['rect']
    svg_path = node.get('asset_svg')
    if svg_path and os.path.isfile(svg_path) and svg_mode == 'shapes':
        try:
            import svg2shapes
            model = svg2shapes.convert(svg_path, r)
            add_svg_shape_group(slide, model, r)
            return None
        except svg2shapes.Unconvertible as e:
            print(f"  svg2shapes fallback ({e}) — svgBlip: {os.path.basename(svg_path)}")
        except Exception as e:
            print(f"  warn: svg2shapes error ({e}) — svgBlip: {os.path.basename(svg_path)}")
    pic = slide.shapes.add_picture(node['asset'], px(r['x']), px(r['y']), px(r['w']), px(r['h']))
    if svg_path and os.path.isfile(svg_path):
        try:
            _attach_svg_blip(slide, pic, svg_path)
        except Exception as e:
            print(f"  warn: svgBlip failed for {os.path.basename(svg_path)} ({e}) — PNG fallback only")
    return pic


def _set_fill(sp, fill):
    if fill is None:
        sp.fill.background()
        return
    if fill['kind'] == 'solid':
        sp.fill.solid()
        sp.fill.fore_color.rgb = rgb(fill['color'])
        a = fill.get('alpha', 1)
        if a < 1:
            el = sp._element.spPr.find(qn('a:solidFill'))
            if el is not None:
                el.find(qn('a:srgbClr')).append(parse_xml(
                    f'<a:alpha xmlns:a="{A_NS}" val="{int(a * 100000)}"/>'))
        return
    # gradient: set a solid first (creates the fill element in schema position), replace
    sp.fill.solid()
    sp.fill.fore_color.rgb = rgb(fill['stops'][0][1])
    solid = sp._element.spPr.find(qn('a:solidFill'))
    stops = ''.join(
        f'<a:gs pos="{int(round(p * 100000))}"><a:srgbClr val="{c.lstrip("#")}">'
        + (f'<a:alpha val="{int(a * 100000)}"/>' if a < 1 else '')
        + '</a:srgbClr></a:gs>'
        for p, c, a in fill['stops'])
    ang = int(round((fill['angle'] % 360) * 60000))
    grad = parse_xml(f'<a:gradFill xmlns:a="{A_NS}" rotWithShape="1"><a:gsLst>{stops}'
                     f'</a:gsLst><a:lin ang="{ang}" scaled="1"/></a:gradFill>')
    solid.getparent().replace(solid, grad)


def _set_stroke(sp, stroke):
    if not stroke:
        sp.line.fill.background()
        return
    sp.line.color.rgb = rgb(stroke['color'])
    sp.line.width = Pt(max(0.25, stroke['width'] * PT_PER_PX))
    ln = sp._element.spPr.find(qn('a:ln'))
    if ln is not None:
        a = stroke.get('alpha', 1)
        if a < 1:
            srgb = ln.find(qn('a:solidFill'))
            if srgb is not None:
                srgb.find(qn('a:srgbClr')).append(parse_xml(
                    f'<a:alpha xmlns:a="{A_NS}" val="{int(a * 100000)}"/>'))
        if stroke.get('cap') == 'rnd':
            ln.set('cap', 'rnd')
        if stroke.get('dash'):
            ln.append(parse_xml(f'<a:prstDash xmlns:a="{A_NS}" val="dash"/>'))


def _set_shadow(sp, shadow):
    if not shadow:
        sp.shadow.inherit = False
        return
    spPr = sp._element.spPr
    dist = math.hypot(shadow['dx'], shadow['dy']) * EMU_PER_PX
    direction = int(round(math.degrees(math.atan2(shadow['dy'], shadow['dx'])) % 360 * 60000))
    spPr.append(parse_xml(
        f'<a:effectLst xmlns:a="{A_NS}"><a:outerShdw blurRad="{int(shadow["blur"] * 2 * EMU_PER_PX)}" '
        f'dist="{int(dist)}" dir="{direction}" rotWithShape="0">'
        f'<a:srgbClr val="{shadow["color"].lstrip("#")}">'
        f'<a:alpha val="{int(shadow["alpha"] * 100000)}"/></a:srgbClr></a:outerShdw></a:effectLst>'))


def _custgeom(sp, n):
    """Swap the placeholder rect's prstGeom for a custGeom carrying the freeform path.
    python-pptx's FreeformBuilder does straight lines only — beziers need raw oxml."""
    bx, by, bw, bh = n['bbox']
    W, H = max(1, int(bw * EMU_PER_PX)), max(1, int(bh * EMU_PER_PX))

    def pt(x, y):
        return (f'<a:pt x="{int((x - bx) * EMU_PER_PX)}" y="{int((y - by) * EMU_PER_PX)}"/>')

    # All subpaths of one shape MUST live inside a single <a:path> element: DrawingML
    # only applies hole-punching (nonzero/even-odd winding) across moveTo/close sequences
    # *within one* <a:path>. Sibling <a:path> elements in a pathLst are each rendered as
    # independently-filled solid regions instead — which is what made letters with counters
    # (e, a, p, ...) render as solid blobs with no holes when each subpath got its own <a:path>.
    cmds = []
    for spath in n['subpaths']:
        cmds.append(f'<a:moveTo>{pt(*spath["start"])}</a:moveTo>')
        for seg in spath['segs']:
            if seg[0] == 'l':
                cmds.append(f'<a:lnTo>{pt(seg[1], seg[2])}</a:lnTo>')
            else:
                cmds.append(f'<a:cubicBezTo>{pt(seg[1], seg[2])}{pt(seg[3], seg[4])}{pt(seg[5], seg[6])}</a:cubicBezTo>')
        if spath['closed']:
            cmds.append('<a:close/>')
    geom = parse_xml(
        f'<a:custGeom xmlns:a="{A_NS}"><a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
        f'<a:rect l="0" t="0" r="{W}" b="{H}"/>'
        f'<a:pathLst><a:path w="{W}" h="{H}">{"".join(cmds)}</a:path></a:pathLst></a:custGeom>')
    spPr = sp._element.spPr
    prst = spPr.find(qn('a:prstGeom'))
    spPr.replace(prst, geom)


def add_svg_shape_group(slide, model, rect):
    """Emit the svg2shapes model as ONE group shape: children keep absolute slide
    coordinates via the chOff=off / chExt=ext trick, so the group renders in place and
    scales/moves as a unit; ungroup in PowerPoint for per-shape editing."""
    spTree = slide.shapes._spTree
    built = []
    for n in model:
        if n['t'] == 'rect':
            shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if n.get('rx') else MSO_SHAPE.RECTANGLE
            sp = slide.shapes.add_shape(shape_type, px(n['x']), px(n['y']),
                                        px(max(n['w'], 0.5)), px(max(n['h'], 0.5)))
            if n.get('rx'):
                try:
                    sp.adjustments[0] = min(0.5, n['rx'] / max(1.0, min(n['w'], n['h'])))
                except Exception:
                    pass
            _set_fill(sp, n['fill']); _set_stroke(sp, n['stroke']); _set_shadow(sp, n['shadow'])
        elif n['t'] == 'ellipse':
            sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(n['x']), px(n['y']),
                                        px(max(n['w'], 0.5)), px(max(n['h'], 0.5)))
            _set_fill(sp, n['fill']); _set_stroke(sp, n['stroke']); _set_shadow(sp, n['shadow'])
        elif n['t'] == 'free':
            bx, by, bw, bh = n['bbox']
            sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(bx), px(by), px(bw), px(bh))
            _custgeom(sp, n)
            _set_fill(sp, n['fill']); _set_stroke(sp, n['stroke']); _set_shadow(sp, n['shadow'])
        elif n['t'] == 'text':
            sp = _svg_textbox(slide, n)
        else:
            continue
        built.append(sp._element)

    if not built:
        return
    # group container
    max_id = max([int(e.get('id')) for e in spTree.iter() if e.tag == qn('p:cNvPr') and e.get('id')] or [1])
    off_x, off_y = int(rect['x'] * EMU_PER_PX), int(rect['y'] * EMU_PER_PX)
    ext_x, ext_y = int(rect['w'] * EMU_PER_PX), int(rect['h'] * EMU_PER_PX)
    grp = parse_xml(
        f'<p:grpSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        f'xmlns:a="{A_NS}">'
        f'<p:nvGrpSpPr><p:cNvPr id="{max_id + 1}" name="Diagram"/>'
        f'<p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
        f'<p:grpSpPr><a:xfrm>'
        f'<a:off x="{off_x}" y="{off_y}"/><a:ext cx="{ext_x}" cy="{ext_y}"/>'
        f'<a:chOff x="{off_x}" y="{off_y}"/><a:chExt cx="{ext_x}" cy="{ext_y}"/>'
        f'</a:xfrm></p:grpSpPr></p:grpSp>')
    for el in built:
        spTree.remove(el)
        grp.append(el)
    spTree.append(grp)


def _svg_textbox(slide, n):
    size = n['size']
    est_w = max(len(n['text']) * size * 0.75, size * 2)
    if n.get('baseline') in ('middle', 'central'):
        top = n['y'] - size * 0.62
    else:
        top = n['y'] - size * 0.85
    if n['anchor'] == 'middle':
        left, align = n['x'] - est_w / 2, PP_ALIGN.CENTER
    elif n['anchor'] == 'end':
        left, align = n['x'] - est_w, PP_ALIGN.RIGHT
    else:
        left, align = n['x'], PP_ALIGN.LEFT
    tb = slide.shapes.add_textbox(px(left), px(top), px(est_w), px(size * 1.4))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    for fit in ('a:normAutofit', 'a:spAutoFit'):
        el = bodyPr.find(qn(fit))
        if el is not None:
            bodyPr.remove(el)
    para = tf.paragraphs[0]
    para.alignment = align
    rn = para.add_run()
    rn.text = n['text']
    rn.font.name = font_for_weight(n['weight'])
    rn.font.size = Pt(round(size * PT_PER_PX, 1))
    rn.font.bold = False
    rn.font.color.rgb = rgb(n['color'])
    if n.get('grad'):
        solid = rn.font._rPr.find(qn('a:solidFill'))
        if solid is not None:
            stops = ''.join(
                f'<a:gs pos="{int(round(p * 100000))}"><a:srgbClr val="{c.lstrip("#")}"/></a:gs>'
                for p, c, a in n['grad']['stops'])
            ang = int(round((n['grad']['angle'] % 360) * 60000))
            solid.getparent().replace(solid, parse_xml(
                f'<a:gradFill xmlns:a="{A_NS}"><a:gsLst>{stops}</a:gsLst>'
                f'<a:lin ang="{ang}" scaled="1"/></a:gradFill>'))
    if n.get('rot'):
        tb.rotation = n['rot']
    return tb


def _attach_svg_blip(slide, pic, svg_path):
    from pptx.opc.package import Part
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    part = slide.part
    pkg = part.package
    with open(svg_path, 'rb') as f:
        blob = f.read()
    partname = pkg.next_partname('/ppt/media/image%d.svg')
    svg_part = Part(partname, 'image/svg+xml', pkg, blob)
    rId = part.relate_to(svg_part, RT.IMAGE)
    blip = pic._element.blipFill.find(qn('a:blip'))
    ext_lst = parse_xml(
        f'<a:extLst xmlns:a="{A_NS}">'
        f'<a:ext uri="{SVG_EXT_URI}">'
        f'<asvg:svgBlip xmlns:asvg="http://schemas.microsoft.com/office/drawing/2016/SVG/main" '
        f'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
        f'r:embed="{rId}"/></a:ext></a:extLst>')
    blip.append(ext_lst)


# Chart series order comes from the pack's `viz` role list (hex, '#'-stripped for
# python-pptx's RGBColor.from_string). A pack that fills no `viz` role gets neutral
# greys, not a colour from some other brand's palette — same reasoning as
# AXIS_LINE_COLOR below.
CHART_COLORS = [c.lstrip('#') for c in _DS.palette('viz')] or ['999999', '6A6A6A', '3D3D3D']
# No literal fallback: a pack that fills no `rule` role should show a neutral grey, not a
# colour from some other brand's palette.
AXIS_LINE_COLOR = (_DS.color('rule') or '#999999').lstrip('#')


def add_chart(slide, node):
    """Native, data-editable PPTX chart from scraped ECharts data (opt-in path)."""
    ch = node['chart']
    type_map = {'bar': XL_CHART_TYPE.COLUMN_CLUSTERED, 'line': XL_CHART_TYPE.LINE,
                'pie': XL_CHART_TYPE.PIE}
    xl_type = type_map.get(ch.get('type'), XL_CHART_TYPE.COLUMN_CLUSTERED)
    data = CategoryChartData()
    data.categories = [str(c) for c in ch['categories']] or [str(i + 1) for i in range(len(ch['series'][0]['data']))]
    for s in ch['series']:
        data.add_series(s['name'] or 'Series', s['data'])
    r = node['rect']
    gframe = slide.shapes.add_chart(xl_type, px(r['x']), px(r['y']), px(r['w']), px(r['h']), data)
    chart = gframe.chart
    chart.has_title = False
    try:
        chart.has_legend = len(ch['series']) > 1
        for i, plot in enumerate(chart.plots):
            for j, series in enumerate(plot.series):
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = RGBColor.from_string(
                    CHART_COLORS[j % len(CHART_COLORS)])
        for axis in (chart.category_axis, chart.value_axis):
            axis.format.line.color.rgb = RGBColor.from_string(AXIS_LINE_COLOR)
            axis.tick_labels.font.size = Pt(10.5)
            axis.tick_labels.font.name = font_for_weight(400)
    except Exception:
        pass  # styling is best-effort; the data is what matters
    return gframe


def set_gradient_background(slide, grad):
    """Native slide-background gradient (covers/dividers stay editable). CSS angle
    (0deg = up, clockwise) maps to DrawingML ang (0 = right, clockwise, 1/60000 deg)."""
    ang = int(round(((grad['angle'] - 90) % 360) * 60000))
    stops = ''.join(
        f'<a:gs pos="{int(round(s["pos"] * 1000))}"><a:srgbClr val="{s["color"].lstrip("#")}"/></a:gs>'
        for s in grad['stops'])
    grad_xml = (f'<a:gradFill xmlns:a="{A_NS}" rotWithShape="1"><a:gsLst>{stops}</a:gsLst>'
                f'<a:lin ang="{ang}" scaled="1"/></a:gradFill>')
    bg_xml = (f'<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
              f'xmlns:a="{A_NS}"><p:bgPr>{grad_xml}<a:effectLst/></p:bgPr></p:bg>')
    cSld = slide._element.find(qn('p:cSld'))
    old = cSld.find(qn('p:bg'))
    if old is not None:
        cSld.remove(old)
    cSld.insert(0, parse_xml(bg_xml))


def build(ir_paths, output, svg_mode='shapes'):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    for path in ir_paths:
        ir = json.load(open(path, encoding='utf-8'))
        slide = prs.slides.add_slide(blank)

        meta = ir['slide']
        if meta.get('bgAsset'):
            # unparseable/layered background — whole rendered slide as base layer;
            # boxes/text above it would duplicate, so these slides carry ONLY the picture
            slide.shapes.add_picture(meta['bgAsset'], 0, 0, SLIDE_W, SLIDE_H)
            continue
        if meta.get('bgGradient'):
            set_gradient_background(slide, meta['bgGradient'])
        else:
            bg = slide.background
            bg.fill.solid()
            bg.fill.fore_color.rgb = rgb(meta.get('bg', '#FFFFFF'))

        for node in ir['nodes']:
            if node['kind'] == 'box':
                add_box(slide, node)
            elif node['kind'] == 'raster':
                add_raster(slide, node, svg_mode=svg_mode)
            elif node['kind'] == 'chart':
                add_chart(slide, node)
            elif node['kind'] == 'text':
                add_text(slide, node)

    prs.save(output)
    return len(ir_paths)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('ir', nargs='+', help='IR JSON files in slide order')
    ap.add_argument('--output', required=True)
    ap.add_argument('--svg-blip', action='store_true',
                    help='embed SVGs as vector pictures instead of native shape groups')
    args = ap.parse_args()
    n = build(args.ir, args.output, svg_mode='blip' if args.svg_blip else 'shapes')
    print(f"{n} slide(s) -> {args.output}")


if __name__ == '__main__':
    main()
